"""Lectores multiformato que convierten archivos en documentos de LangChain.

Cada lector conserva metadatos suficientes para rastrear una respuesta hasta el
archivo, página, hoja o diapositiva de origen. Los archivos se procesan localmente.
"""

from __future__ import annotations

from collections.abc import Iterator
from pathlib import Path

import pandas as pd
from docx import Document as WordDocument
from langchain_core.documents import Document
from pptx import Presentation
from pypdf import PdfReader

SUPPORTED_EXTENSIONS = frozenset(
    {".pdf", ".docx", ".xlsx", ".xlsm", ".xls", ".csv", ".pptx", ".txt", ".md"}
)


class DocumentLoadError(RuntimeError):
    """Explica qué archivo no pudo convertirse a texto."""


class MultiFormatDocumentLoader:
    """Carga PDF, Word, Excel, PowerPoint, CSV y archivos de texto."""

    __slots__ = ("directory",)

    def __init__(self, directory: Path) -> None:
        """Valida y conserva la carpeta que contiene la base de información."""

        self.directory = directory.resolve()
        if not self.directory.is_dir():
            raise DocumentLoadError(
                f"No existe la carpeta de información: {self.directory}"
            )

    def source_files(self) -> tuple[Path, ...]:
        """Devuelve archivos compatibles, ordenados y sin temporales de Office."""

        return tuple(
            sorted(
                (
                    path
                    for path in self.directory.iterdir()
                    if path.is_file()
                    and not path.name.startswith("~$")
                    and path.suffix.casefold() in SUPPORTED_EXTENSIONS
                ),
                key=lambda path: path.name.casefold(),
            )
        )

    def lazy_load(self) -> Iterator[Document]:
        """Procesa cada archivo solamente cuando el consumidor solicita otro."""

        files = self.source_files()
        if not files:
            raise DocumentLoadError(
                f"No hay documentos compatibles en {self.directory}."
            )
        for path in files:
            try:
                yield from _load_file(path)
            except DocumentLoadError:
                raise
            except Exception as exc:
                raise DocumentLoadError(
                    f"No se pudo leer {path.name}: {exc}"
                ) from exc

    def load(self) -> list[Document]:
        """Carga todos los documentos; resulta práctico para bases pequeñas."""

        return list(self.lazy_load())


def _load_file(path: Path) -> Iterator[Document]:
    """Selecciona el lector apropiado usando la extensión del archivo."""

    extension = path.suffix.casefold()
    if extension == ".pdf":
        yield from _load_pdf(path)
    elif extension == ".docx":
        yield from _load_word(path)
    elif extension in {".xlsx", ".xlsm", ".xls"}:
        yield from _load_excel(path)
    elif extension == ".csv":
        yield from _load_csv(path)
    elif extension == ".pptx":
        yield from _load_powerpoint(path)
    elif extension in {".txt", ".md"}:
        yield _load_text(path)
    else:
        raise DocumentLoadError(f"Formato no compatible: {path.name}")


def _base_metadata(path: Path, document_type: str) -> dict[str, str]:
    """Crea metadatos comunes sin incluir contenido ni información sensible."""

    return {
        "source": str(path.resolve()),
        "filename": path.name,
        "document_type": document_type,
    }


def _load_pdf(path: Path) -> Iterator[Document]:
    """Extrae cada página PDF mediante pypdf."""

    reader = PdfReader(path)
    metadata = _base_metadata(path, "pdf")
    for page_number, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            yield Document(
                page_content=text,
                metadata={**metadata, "page": page_number},
            )


def _load_word(path: Path) -> Iterator[Document]:
    """Extrae párrafos y tablas de un documento Word mediante python-docx."""

    word_document = WordDocument(path)
    parts = [
        paragraph.text.strip()
        for paragraph in word_document.paragraphs
        if paragraph.text.strip()
    ]
    for table in word_document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    text = "\n".join(parts).strip()
    if text:
        yield Document(
            page_content=text,
            metadata=_base_metadata(path, "word"),
        )


def _load_excel(path: Path) -> Iterator[Document]:
    """Convierte cada hoja de Excel en texto tabular usando pandas."""

    sheets = pd.read_excel(path, sheet_name=None, dtype=str)
    metadata = _base_metadata(path, "excel")
    for sheet_name, frame in sheets.items():
        clean_frame = frame.fillna("")
        text = clean_frame.to_csv(index=False).strip()
        if text:
            yield Document(
                page_content=text,
                metadata={**metadata, "sheet": str(sheet_name)},
            )


def _load_csv(path: Path) -> Iterator[Document]:
    """Lee CSV con pandas y conserva todas sus columnas como texto."""

    frame = pd.read_csv(path, dtype=str).fillna("")
    text = frame.to_csv(index=False).strip()
    if text:
        yield Document(
            page_content=text,
            metadata=_base_metadata(path, "csv"),
        )


def _load_powerpoint(path: Path) -> Iterator[Document]:
    """Extrae cuadros de texto y tablas de cada diapositiva PowerPoint."""

    presentation = Presentation(path)
    metadata = _base_metadata(path, "powerpoint")
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        content = "\n".join(parts).strip()
        if content:
            yield Document(
                page_content=content,
                metadata={**metadata, "slide": slide_number},
            )


def _load_text(path: Path) -> Document:
    """Lee archivos TXT o Markdown usando UTF-8 con marca BOM opcional."""

    return Document(
        page_content=path.read_text(encoding="utf-8-sig").strip(),
        metadata=_base_metadata(path, path.suffix.casefold().lstrip(".")),
    )

