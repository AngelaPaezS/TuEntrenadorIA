"""Pruebas de los lectores para los formatos solicitados."""

from pathlib import Path
import unittest
from unittest.mock import patch

import pandas as pd
from docx import Document as WordDocument
from pptx import Presentation
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

from tu_entrenador_ia.document_loaders import MultiFormatDocumentLoader
from tu_entrenador_ia.retrieval import DocumentCorpus

PROJECT_ROOT = Path(__file__).resolve().parents[1]


class MultiFormatDocumentLoaderTests(unittest.TestCase):
    """Crea archivos mínimos y comprueba texto y metadatos."""

    def setUp(self) -> None:
        """Prepara una carpeta aislada dentro del proyecto."""

        self.directory = PROJECT_ROOT / ".tmp" / "document_loader_tests"
        self.directory.mkdir(parents=True, exist_ok=True)
        _create_word(self.directory / "ejemplo.docx")
        _create_excel(self.directory / "ejemplo.xlsx")
        _create_powerpoint(self.directory / "ejemplo.pptx")
        _create_pdf(self.directory / "ejemplo.pdf")
        (self.directory / "ejemplo.csv").write_text(
            "ejercicio,series\nSentadilla,3\n",
            encoding="utf-8",
        )
        (self.directory / "ejemplo.txt").write_text(
            "Documento de seguridad para principiantes.",
            encoding="utf-8",
        )

    def test_loads_all_supported_office_formats(self) -> None:
        """Cada formato debe producir contenido y metadatos trazables."""

        documents = MultiFormatDocumentLoader(self.directory).load()
        document_types = {
            document.metadata["document_type"] for document in documents
        }
        self.assertEqual(
            {"word", "excel", "powerpoint", "pdf", "csv", "txt"},
            document_types,
        )
        self.assertTrue(all(document.page_content for document in documents))
        self.assertTrue(all("filename" in document.metadata for document in documents))

    def test_excel_and_powerpoint_keep_locations(self) -> None:
        """Hojas y diapositivas deben aparecer en los metadatos."""

        documents = MultiFormatDocumentLoader(self.directory).load()
        excel_document = next(
            item for item in documents if item.metadata["document_type"] == "excel"
        )
        slide_document = next(
            item
            for item in documents
            if item.metadata["document_type"] == "powerpoint"
        )
        self.assertEqual("Rutina", excel_document.metadata["sheet"])
        self.assertEqual(1, slide_document.metadata["slide"])

    def test_corpus_uses_valid_cache(self) -> None:
        """Una segunda carga no debe reprocesar archivos sin cambios."""

        cache_path = self.directory / ".cache" / "chunks.json"
        first_corpus = DocumentCorpus.build(self.directory, cache_path)
        self.assertTrue(cache_path.is_file())
        with patch.object(
            MultiFormatDocumentLoader,
            "load",
            side_effect=AssertionError("No debe recargar documentos"),
        ):
            second_corpus = DocumentCorpus.build(self.directory, cache_path)
        self.assertEqual(
            len(first_corpus.documents),
            len(second_corpus.documents),
        )

    def test_search_finds_text_from_excel(self) -> None:
        """El índice común debe consultar contenido proveniente de Excel."""

        corpus = DocumentCorpus.build(self.directory)
        results = corpus.search("Puente glúteos doce repeticiones", limit=3)
        self.assertTrue(results)
        self.assertTrue(
            any(
                result.document.metadata["document_type"] == "excel"
                for result in results
            )
        )


def _create_word(path: Path) -> None:
    """Genera un Word pequeño con párrafo y tabla."""

    document = WordDocument()
    document.add_heading("Manual", level=1)
    document.add_paragraph("Ejercicios autorizados para principiantes.")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Ejercicio"
    table.cell(0, 1).text = "Sentadilla"
    document.save(path)


def _create_excel(path: Path) -> None:
    """Genera un libro con una hoja nombrada y datos de ejercicio."""

    frame = pd.DataFrame(
        [{"ejercicio": "Puente de glúteos", "repeticiones": "12"}]
    )
    frame.to_excel(path, sheet_name="Rutina", index=False)


def _create_powerpoint(path: Path) -> None:
    """Genera una presentación con una diapositiva de texto."""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Seguridad"
    slide.placeholders[1].text = "Detenerse si aparece dolor intenso."
    presentation.save(path)


def _create_pdf(path: Path) -> None:
    """Genera un PDF mínimo con texto extraíble sin dependencias adicionales."""

    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_reference = writer._add_object(font)
    page[NameObject("/Resources")] = DictionaryObject(
        {
            NameObject("/Font"): DictionaryObject(
                {NameObject("/F1"): font_reference}
            )
        }
    )
    stream = DecodedStreamObject()
    stream.set_data(
        b"BT /F1 12 Tf 72 720 Td (Rutina PDF para principiantes) Tj ET"
    )
    page[NameObject("/Contents")] = writer._add_object(stream)
    with path.open("wb") as output:
        writer.write(output)


if __name__ == "__main__":
    unittest.main()
